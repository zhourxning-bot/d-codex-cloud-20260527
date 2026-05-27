import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x4F, 0xAC, 0xFF)
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x43)
ACCENT_RED = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT_GREEN = RGBColor(0x4E, 0xCB, 0x71)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_CARD = RGBColor(0x25, 0x25, 0x3D)
TABLE_HEADER_BG = RGBColor(0x33, 0x33, 0x50)
TABLE_ROW_ALT = RGBColor(0x2A, 0x2A, 0x40)
TABLE_BORDER = RGBColor(0x44, 0x44, 0x66)


def add_dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_slide_number(slide, num, total=8):
    left = prs.slide_width - Inches(1.0)
    top = prs.slide_height - Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, Inches(0.8), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_title_bar(slide, title_text, subtitle_text=None):
    # Top accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.95), Inches(11.7), Inches(0.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = LIGHT_GRAY


def add_body_text(slide, left, top, width, height, lines, font_size=14, line_spacing=1.4):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, is_bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = is_bold
        p.font.color.rgb = color if color else WHITE
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_CARD
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x50)
    shape.line.width = Pt(1)
    return shape


def add_section_table(slide, left, top, width, rows_data, col_widths, font_size=10):
    """Add a formatted table. rows_data[0] is header."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_height = n_rows * 0.35

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(width), Inches(table_height)
    )
    table = table_shape.table

    # Set column widths
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)

    for ri, row_data in enumerate(rows_data):
        for ci, cell_text in enumerate(row_data):
            cell = table.cell(ri, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = str(cell_text)
            p.font.size = Pt(font_size)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.LEFT if ci >= 2 else PP_ALIGN.CENTER

            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
                p.font.bold = True
                p.font.size = Pt(font_size + 1)
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_CARD

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # thin margins
            cell.margin_left = Pt(3)
            cell.margin_right = Pt(3)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)

    return table_shape


# ============================================================
# SLIDE 1: Title
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_dark_bg(slide1)

# Large accent shape
shape = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.12)
)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

# Main title
txBox = slide1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.9), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "特看台湾 UGC/广告素材核查报告"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE

p2 = tf.add_paragraph()
p2.text = "TikTok Advertising Material Audit — Taiwan Region"
p2.font.size = Pt(20)
p2.font.color.rgb = ACCENT_BLUE
p2.space_before = Pt(8)

# Subtitle / metadata
txBox2 = slide1.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(10.9), Inches(1.2))
tf2 = txBox2.text_frame
tf2.word_wrap = True
meta_lines = [
    ("查询周期：2026-05-16 至 2026-05-25（近10天）", False, LIGHT_GRAY),
    ("数据平台：特看 TabCut", False, LIGHT_GRAY),
    ("生成日期：2026-05-26", False, LIGHT_GRAY),
    ("生成工具：Claude Code（python-pptx）", False, LIGHT_GRAY),
]
for i, (text, bold, color) in enumerate(meta_lines):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = tf2.add_paragraph()
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_after = Pt(4)

add_slide_number(slide1, 1)
add_card(slide1, 0.8, 5.8, 11.7, 0.9)
add_body_text(slide1, 1.0, 5.95, 11.3, 0.7, [
    ("⚠ 重要提示：本报告数据来源为特看广告搜索页可见素材，未包含转化率等非公开指标。"
     "素材排序基于展现量/热度，不等同于已验证的转化率排名。本报告仅供素材参考与方向性分析使用。",
     False, ACCENT_ORANGE)
], font_size=13)

# ============================================================
# SLIDE 2: 查询范围与目的
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide2)
add_title_bar(slide2, "查询范围与目的", "Scope & Objective")
add_slide_number(slide2, 2)

add_card(slide2, 0.8, 1.7, 11.7, 2.0)
add_body_text(slide2, 1.1, 1.85, 11.1, 1.7, [
    ("🎯 查询目标", True, ACCENT_BLUE),
    ("", False, WHITE),
    ("检索台湾地区（region=TW）最近 10 天内发布的 UGC 视频素材，按点击/转化表现排序，取前 20 条。", False, WHITE),
    ("目的：了解当前台湾地区 TikTok 广告/带货视频的内容趋势、达人分布和互动表现。", False, WHITE),
], font_size=15)

add_card(slide2, 0.8, 3.9, 5.7, 2.8)
add_body_text(slide2, 1.1, 4.05, 5.1, 2.5, [
    ("📋 检索条件", True, ACCENT_BLUE),
    ("", False, WHITE),
    ("• 地区：台湾 (region=TW)", False, WHITE),
    ("• 时间范围：2026-05-16 ~ 2026-05-25", False, WHITE),
    ("• 素材类型：UGC 视频", False, WHITE),
    ("• 排序字段：点击/转化优先", False, WHITE),
    ("• 期望条数：前 20 条", False, WHITE),
    ("• 平台：特看 TabCut", False, WHITE),
], font_size=14)

add_card(slide2, 6.8, 3.9, 5.7, 2.8)
add_body_text(slide2, 7.1, 4.05, 5.1, 2.5, [
    ("🔗 数据源链接", True, ACCENT_BLUE),
    ("", False, WHITE),
    ("视频库检索：", False, LIGHT_GRAY),
    ("tabcut.com/.../video-search?region=TW&sortField=video_sold_count...", False, WHITE),
    ("", False, WHITE),
    ("广告素材检索：", False, LIGHT_GRAY),
    ("tabcut.com/.../advertising-search?region=TW&sortField=143...", False, WHITE),
], font_size=12)

# ============================================================
# SLIDE 3: 视频库检索 — No Data 结论
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide3)
add_title_bar(slide3, "视频库检索结论", "Video Search Result")
add_slide_number(slide3, 3)

add_card(slide3, 0.8, 1.7, 11.7, 5.0)
add_body_text(slide3, 1.1, 1.85, 11.1, 4.6, [
    ("🔴 检索结果：No Data", True, ACCENT_RED),
    ("", False, WHITE),
    ("在以下条件下，特看「发现视频」模块返回 No data：", False, WHITE),
    ("", False, WHITE),
    ("• 地区：台湾 (TW)", False, WHITE),
    ("• 时间：2026-05-16 至 2026-05-25", False, WHITE),
    ("• 排序：按转化/点击表现排序", False, WHITE),
    ("• 视频类型：带货视频 (itemVideoFlag=1)", False, WHITE),
    ("", False, WHITE),
    ("🟡 含义解读", True, ACCENT_ORANGE),
    ("", False, WHITE),
    ("1. 特看视频库在台湾地区的近期 UGC 带货视频数据覆盖有限，", False, WHITE),
    ("   可能因平台数据采集策略、区域覆盖或视频识别规则导致。", False, WHITE),
    ("2. 这并不意味着台湾地区没有 TikTok 带货视频，", False, WHITE),
    ("   而是特看当前该模块下无可用结构化数据。", False, WHITE),
    ("3. 因此本报告转向广告搜索页（advertising-search）获取可参考的素材信息。", False, WHITE),
], font_size=14)

# ============================================================
# SLIDE 4: 广告素材列表 (1-10)
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide4)
add_title_bar(slide4, "可见 Country.TW 广告素材列表 (1-10/20)", "Advertising Material List — Page 1")
add_slide_number(slide4, 4)

headers = ['#', '达人/账号', '素材标题', '投放时间', '展现量', '热度', '点赞比']
data = [
    headers,
    ['1', 'OSCALofficial', 'Oscal TIGER 12: Official Introduction...', '2024-09-11 ~ 10-15', '28.4M', '1.5M', '5.28%'],
    ['2', 'OSCALofficial', 'Oscal Pad 18: 11 Inches. Feel More...', '2024-09-10 ~ 10-15', '7.2M', '486.4K', '6.68%'],
    ['3', '凤凰卫视PhoenixTV', 'Paris2024 Olympics — Russian Coach...', '2024-09-09 ~ 09-14', '3M', '206.9K', '6.75%'],
    ['4', '凤凰卫视PhoenixTV', 'KIDS on the First Day of Kindergarten…', '2024-09-13 ~ 09-14', '1.6M', '64K', '3.55%'],
    ['5', '凤凰卫视PhoenixTV', 'Massive Lightning Strikes Hainan...', '2024-09-09 ~ 09-12', '1.4M', '59.9K', '3.88%'],
    ['6', '安哥碎碎念', '安哥好车推荐—莲花LOTUS超跑…', '2024-09-11 ~ 09-11', '475K', '50.7K', '10.48%'],
    ['7', '(未显示)', 'Browse and download videos privately.', '2025-09-27 ~ 09-27', '2.1M', '39.3K', '1.89%'],
    ['8', '安哥碎碎念', '开3000万豪车吃60块炒饭！Feat…', '2024-09-09 ~ 09-14', '606.7K', '33.7K', '5.18%'],
    ['9', '辰馨成就懒女人', '你一定要抓住 #开店 #药局 #拍照…', '2024-10-12 ~ 10-12', '2.1M', '30.4K', '1.23%'],
    ['10', '(未显示)', '其实我天生社恐 #社恐 #自媒体…', '2025-08-18 ~ 08-18', '627.4K', '29K', '4.54%'],
]

col_widths = [0.4, 1.8, 5.0, 1.6, 1.0, 0.8, 0.8]
add_section_table(slide4, 0.8, 1.8, 11.4, data, col_widths, font_size=10)

add_body_text(slide4, 0.8, 6.2, 11.7, 0.8, [
    ("注：排序依据特看广告搜索页默认热度排序（sortField=143）。「点赞比」= 热度/展现量，作为点击/互动质量参考。", False, LIGHT_GRAY),
], font_size=11)

# ============================================================
# SLIDE 5: 广告素材列表 (11-20)
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide5)
add_title_bar(slide5, "可见 Country.TW 广告素材列表 (11-20/20)", "Advertising Material List — Page 2")
add_slide_number(slide5, 5)

data2 = [
    headers,
    ['11', 'Jay & Sharon', 'We promised to keep our own gifts... @PayPal', '2025-12-13 ~ 12-21', '6.8M', '27.2K', '0.38%'],
    ['12', '(未显示)', 'Have you ever bought perfume this cheap?...', '2026-01-04 ~ 01-04', '4.8M', '24.8K', '0.50%'],
    ['13', '(未显示)', '不是谁都能接班，接得住责任… #澄霖生医', '2025-08-18 ~ 08-18', '566K', '21K', '3.66%'],
    ['14', '凤凰卫视PhoenixTV', 'You never expected what would happen…😨', '2024-09-09 ~ 09-09', '393.6K', '20.2K', '4.99%'],
    ['15', '凤凰卫视PhoenixTV', 'After TYPHOON YAGI comes harvest...', '2024-09-13 ~ 09-14', '267.4K', '15.6K', '5.70%'],
    ['16', '(未显示)', '勇闯村庄遇险 收获丰富翡翠美镯…', '2025-09-12 ~ 09-12', '1.5M', '13.6K', '0.86%'],
    ['17', 'foodiesushiqueen', '#Ad NEW Wood Fired-Style Pizzas... @DiGiorno', '2025-10-03 ~ 10-08', '6M', '11.9K', '0.19%'],
    ['18', 'MSI', 'Prebuilt or custom desktop. Which team?...', '2026-03-29 ~ 03-30', '341.1K', '8.7K', '2.53%'],
    ['19', '(国际通报骗子)', '"国际通报"骗子账号！', '2025-09-13 ~ 09-13', '887.5K', '7.9K', '0.77%'],
    ['20', '(国际通报骗子)', '"国际通报"骗子账号！', '2025-09-12 ~ 09-12', '866.2K', '7.7K', '0.79%'],
]

add_section_table(slide5, 0.8, 1.8, 11.4, data2, col_widths, font_size=10)

add_body_text(slide5, 0.8, 6.2, 11.7, 0.8, [
    ("注：部分达人/账号信息未直接显示。数据日期跨越较大（2024-09 至 2026-03），非全部为近10天新素材。", False, LIGHT_GRAY),
], font_size=11)

# ============================================================
# SLIDE 6: 关键信号分析
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide6)
add_title_bar(slide6, "关键信号分析", "Key Signals")
add_slide_number(slide6, 6)

# Card 1: top-left
add_card(slide6, 0.8, 1.7, 5.7, 2.3)
add_body_text(slide6, 1.1, 1.85, 5.1, 2.0, [
    ("🏆 高展现量头部素材", True, ACCENT_BLUE),
    ("", False, WHITE),
    ("• OSCALofficial — 展现量最高 (28.4M), 两款产品合计 35.6M", False, WHITE),
    ("• Jay & Sharon — 6.8M 展现 (PayPal 合作)", False, WHITE),
    ("• foodiesushiqueen — 6M 展现 (DiGiorno 广告)", False, WHITE),
    ("• 凤凰卫视 — 5条素材入榜, 合计展现 ~6.7M", False, WHITE),
    ("→ 3C 品牌 + 媒体账号占据头部流量", False, ACCENT_ORANGE),
], font_size=13)

# Card 2: top-right
add_card(slide6, 6.8, 1.7, 5.7, 2.3)
add_body_text(slide6, 7.1, 1.85, 5.1, 2.0, [
    ("❤️ 高点赞比（互动质量）", True, ACCENT_GREEN),
    ("", False, WHITE),
    ("• 安哥碎碎念 (Lotus超跑) — 10.48% (最高)", False, WHITE),
    ("• 凤凰卫视 (Paris Olympics) — 6.75%", False, WHITE),
    ("• OSCALofficial (Pad 18) — 6.68%", False, WHITE),
    ("• 凤凰卫视 (Typhoon Yagi) — 5.70%", False, WHITE),
    ("→ 超跑/3C/新闻热点类内容互动率最高", False, ACCENT_ORANGE),
], font_size=13)

# Card 3: bottom-left
add_card(slide6, 0.8, 4.2, 5.7, 2.3)
add_body_text(slide6, 1.1, 4.35, 5.1, 2.0, [
    ("📊 达人分布", True, ACCENT_BLUE),
    ("", False, WHITE),
    ("• 凤凰卫视PhoenixTV — 5条 (25%), 媒体账号高频投放", False, WHITE),
    ("• OSCALofficial — 2条, 3C品牌官方投放", False, WHITE),
    ("• 安哥碎碎念 — 2条, 汽车类KOL", False, WHITE),
    ("• 其他达人各1条, 账号分散度较高", False, WHITE),
    ("• 7条素材未显示达人/账号名称", False, LIGHT_GRAY),
], font_size=13)

# Card 4: bottom-right
add_card(slide6, 6.8, 4.2, 5.7, 2.3)
add_body_text(slide6, 7.1, 4.35, 5.1, 2.0, [
    ("📅 时间分布特征", True, ACCENT_GREEN),
    ("", False, WHITE),
    ("• 2024-09：10条 (50%), 1年前旧素材仍占主导", False, WHITE),
    ("• 2025-08~10：5条, 包括DiGiorno等品牌投放", False, WHITE),
    ("• 2025-12~2026-03：4条, PayPal/MSI等品牌合作", False, WHITE),
    ("• 近10天（2026-05-16~25）：0条新素材", False, ACCENT_RED),
    ("→ 搜索结果以历史素材为主，近期活跃度低", False, ACCENT_ORANGE),
], font_size=13)

# ============================================================
# SLIDE 7: 局限性
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide7)
add_title_bar(slide7, "数据局限性", "Limitations & Caveats")
add_slide_number(slide7, 7)

add_card(slide7, 0.8, 1.7, 11.7, 5.0)
add_body_text(slide7, 1.1, 1.85, 11.1, 4.6, [
    ("🔶 方法局限性", True, ACCENT_ORANGE),
    ("", False, WHITE),
    ("1. 不可用字段：转化率和独立视频详情链接未在页面 HTML 中暴露，无法获取核心转化数据。", False, WHITE),
    ("2. 「点赞比」仅为互动质量的方向性代理指标，不等同于转化率或 ROI。", False, WHITE),
    ("3. 展现量和热度为平台聚合数据，未经独立验证。", False, WHITE),
    ("", False, WHITE),
    ("🔶 数据局限性", True, ACCENT_ORANGE),
    ("", False, WHITE),
    ("4. 视频库模块（video-search）对台湾地区 + 近10天 + UGC 带货视频返回 No data，", False, WHITE),
    ("   无法进行原计划的 UGC 转化排名分析。", False, WHITE),
    ("5. 广告搜索页返回的前 20 条并非全部为近10天新素材；大部分投放时间为 2024 年 9 月，", False, WHITE),
    ("   时间覆盖跨度超过 1.5 年，时效性不可控。", False, WHITE),
    ("6. 7/20 条素材未显示达人/账号名称，来源信息不完整。", False, WHITE),
    ("", False, WHITE),
    ("🔶 解读警示", True, ACCENT_RED),
    ("", False, WHITE),
    ("本报告素材清单为特看广告搜索页的可见「广告素材参考列表」，", False, WHITE),
    ("并非经验证的转化率排名。任何以本清单为基础的投放决策应结合独立验证数据。", False, ACCENT_ORANGE),
], font_size=13)

# ============================================================
# SLIDE 8: 后续验证步骤
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide8)
add_title_bar(slide8, "后续验证步骤", "Next Verification Steps")
add_slide_number(slide8, 8)

steps = [
    ("Step 1", "扩大检索范围", "尝试调整特看视频库检索条件：放宽时间跨度、切换排序字段（互动量/点赞数）、调整视频类型过滤，确认是否能获取台湾地区 UGC 带货视频数据。"),
    ("Step 2", "跨平台交叉验证", "通过 TikTok 广告库、第三方监测工具（如 Kalodata、Tabcut 其他模块）对广告搜索页返回的高展现素材进行转化率交叉验证。"),
    ("Step 3", "达人深度分析", "针对本报告识别的高互动达人（安哥碎碎念、凤凰卫视、OSCALofficial），拉取其近 30 天发布内容和互动趋势，评估合作价值。"),
    ("Step 4", "定期复查询", "每周执行一次相同的视频库查询，监测 No data 状态是否变化；当数据可用时立即补充 UGC 转化排名分析。"),
    ("Step 5", "补充转化数据", "若后续特看或其他平台开放转化率字段，重新生成本报告并将「点赞比」替换为实际转化指标。"),
]

y = 1.9
for step_num, step_title, step_desc in steps:
    add_card(slide8, 0.8, y, 11.7, 0.95)
    # Step number badge
    txBox = slide8.shapes.add_textbox(Inches(1.1), Inches(y + 0.15), Inches(1.0), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = step_num
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    add_body_text(slide8, 2.2, y + 0.1, 10.0, 0.8, [
        (step_title, True, WHITE),
        (step_desc, False, LIGHT_GRAY),
    ], font_size=12)
    y += 1.05

# ============================================================
# Save
# ============================================================
output_path = r"D:\跨境素材\特看_台湾UGC视频_近10天核查_ClaudeCode生成.pptx"
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
print(f"Slide count: {len(prs.slides)}")
